from pptx import Presentation
from pptx.util import Pt


class PowerPointGenerator:
    def __init__(self, title_font_size=16, bullet_point_font_size=12):
        self.title_font_size = title_font_size
        self.bullet_point_font_size = bullet_point_font_size
        self.prs = Presentation()

    def add_title(self, topic_title, slide):
        title = slide.shapes.title
        title.text = topic_title
        title_text_frame = title.text_frame
        title_text_frame.paragraphs[0].font.size = Pt(self.title_font_size)

    def add_bullet_points(self, points_array, slide):
        content_placeholder = slide.placeholders[1]  # Index 1 is the content placeholder

        content_text_frame = content_placeholder.text_frame

        for point in points_array:
            p = content_text_frame.add_paragraph()

            p.text = point.removesuffix('\n')

        for paragraph in range(len(content_text_frame.paragraphs)):
            content_text_frame.paragraphs[paragraph].font.size = Pt(self.bullet_point_font_size)

    def add_slide(self):
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)

        return slide

    def save_pptx(self, path):
        self.prs.save(path)
